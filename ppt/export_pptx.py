"""Export the SpringBoot Takeout Platform PPT to PPTX format."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme colors ──
INK = RGBColor(0x0A, 0x0A, 0x0B)       # #0a0a0b
PAPER = RGBColor(0xF1, 0xEF, 0xEA)     # #f1efea
PAPER_TINT = RGBColor(0xE8, 0xE5, 0xDE)  # #e8e5de
INK_TINT = RGBColor(0x18, 0x18, 0x1A)  # #18181a
MUTED = RGBColor(0x80, 0x80, 0x80)     # for secondary text
ACCENT = RGBColor(0x3A, 0x3A, 0x3C)    # subtle accent

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_para(tf, text, font_size=18, color=INK, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Microsoft YaHei UI", spacing_before=0, spacing_after=0):
    """Add a paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(spacing_after)
    return p


def add_mixed_para(tf, segments, alignment=PP_ALIGN.LEFT, spacing_before=0, spacing_after=0):
    """Add a paragraph with mixed formatting segments.
    Each segment: (text, font_size, color, bold, font_name)
    """
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(spacing_before)
    p.space_after = Pt(spacing_after)

    for i, (text, font_size, color, bold, font_name) in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return p


def add_card_box(slide, left, top, width, height, border_color=None):
    """Add a rounded rectangle card shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF6, 0xF2) if slide.background.fill.fore_color.rgb == PAPER else RGBColor(0x1A, 0x1A, 0x1C)
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def create_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════
    # Page 1: Cover (hero dark)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK)

    # Chrome top-left
    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "外卖点餐平台 · 微服务架构", font_size=9, color=MUTED,
             font_name="Consolas", spacing_after=0)

    # Main content area
    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(10), Inches(5))

    add_para(tf, "SPRING CLOUD MICROSERVICES", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=16)

    add_para(tf, "外卖点餐平台", font_size=54, color=PAPER, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=8)

    add_para(tf, "微服务架构设计与实现", font_size=24, color=PAPER_TINT,
             font_name="Microsoft YaHei UI", spacing_after=24)

    add_para(tf, "基于 Spring Boot 3.3+ 构建的校园级外卖点餐系统，强调高并发、实时性、分布式一致性、智能调度与全链路可观测。",
             font_size=16, color=PAPER_TINT, spacing_after=32)

    # Meta row
    add_mixed_para(tf, [
        ("JDK 21", 11, MUTED, False, "Consolas"),
        ("  ·  ", 11, MUTED, False, "Consolas"),
        ("Spring Cloud 2025", 11, MUTED, False, "Consolas"),
        ("  ·  ", 11, MUTED, False, "Consolas"),
        ("8 Microservices", 11, MUTED, False, "Consolas"),
    ])

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "一场关于云原生 · 微服务 · 事件驱动的技术分享", font_size=9, color=MUTED,
             font_name="Consolas")

    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "— 2026 —", font_size=9, color=MUTED, font_name="Consolas",
             alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 2: Overview (light)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PAPER)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "项目总览 · Overview", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "01 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(4.5))

    add_para(tf, "OVERVIEW · 总览", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=12)

    add_para(tf, "接近美团校园版", font_size=42, color=INK, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=12)

    add_para(tf, "一个完整的外卖点餐微服务平台，涵盖用户、商家、订单、支付、配送、搜索、通知全链路。",
             font_size=16, color=INK_TINT, spacing_after=36)

    # Stat cards
    card_data = [
        ("JDK", "21", "最新 LTS，虚拟线程支持"),
        ("Microservices", "8", "核心服务拆分"),
        ("Spring Cloud", "2025", "Northfields 最新版"),
    ]
    card_w = Inches(3.2)
    card_h = Inches(2.0)
    card_y = Inches(3.8)
    gap = Inches(0.6)
    start_x = Inches(0.8)

    for i, (label, number, note) in enumerate(card_data):
        x = start_x + i * (card_w + gap)
        # Card background
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PAPER_TINT
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.3), card_y + Inches(0.2), card_w - Inches(0.6), card_h - Inches(0.3))
        add_para(tf, label.upper(), font_size=9, color=MUTED, font_name="Consolas", spacing_after=4)
        add_para(tf, number, font_size=42, color=INK, bold=True,
                 font_name="Playfair Display", spacing_after=4)
        add_para(tf, note, font_size=12, color=INK_TINT, spacing_after=0)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "项目 · 外卖点餐平台", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "01 · Overview", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 3: Tech Stack I (dark)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "技术栈 · Tech Stack", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "02 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5))

    add_para(tf, "PART I · 基础设施层", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "核心技术栈", font_size=42, color=PAPER, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=28)

    rows = [
        ("微服务治理", "Spring Cloud Alibaba + Nacos 3.1.x（3 节点集群，服务注册发现 + 配置中心）"),
        ("API 网关", "Spring Cloud Gateway（统一鉴权、限流、日志、跨域、路由）"),
        ("ORM", "MyBatis-Plus（代码生成、逻辑删除、分页插件）"),
        ("缓存", "Redis 7.x（分布式锁、购物车、Session、实时数据、ZSet 抢单）"),
        ("消息队列", "RabbitMQ（可靠异步通知、支付回调）+ Kafka（高吞吐订单事件流、骑手位置上报）"),
    ]

    for key, val in rows:
        # Separator line
        add_para(tf, "─" * 80, font_size=6, color=RGBColor(0x40, 0x40, 0x40),
                 font_name="Consolas", spacing_before=8, spacing_after=2)
        add_mixed_para(tf, [
            (f"{key}", 15, PAPER, True, "Microsoft YaHei UI"),
            (f"    {val}", 14, PAPER_TINT, False, "Microsoft YaHei UI"),
        ], spacing_after=4)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "技术栈 · 基础设施层", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "02 · Tech Stack I", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 4: Tech Stack II (light)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PAPER)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "技术栈 · Tech Stack", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "03 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(5.5))

    add_para(tf, "PART II · 业务能力层", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "核心技术栈", font_size=42, color=INK, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=28)

    rows = [
        ("搜索引擎", "Elasticsearch 8.x（商家、菜品、订单全文搜索与聚合）"),
        ("分布式事务", "Seata 2.5+（AT / Saga 模式，下单扣库存、支付等一致性保障）"),
        ("限流熔断", "Sentinel（流量防护、降级、热点参数限流）"),
        ("实时通信", "WebSocket + MQTT（骑手位置上报、订单状态推送）"),
        ("地理位置", "高德 / 腾讯地图 API（距离计算、ETA、地理围栏）"),
        ("监控观测", "Prometheus + Grafana + SkyWalking（全链路追踪 + 指标面板）"),
    ]

    for key, val in rows:
        add_para(tf, "─" * 80, font_size=6, color=PAPER_TINT,
                 font_name="Consolas", spacing_before=8, spacing_after=2)
        add_mixed_para(tf, [
            (f"{key}", 15, INK, True, "Microsoft YaHei UI"),
            (f"    {val}", 14, INK_TINT, False, "Microsoft YaHei UI"),
        ], spacing_after=4)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "技术栈 · 业务能力层", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "03 · Tech Stack II", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 5: Architecture (dark)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "架构 · Architecture", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "04 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(2))
    add_para(tf, "ARCHITECTURE · 架构设计", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "六大架构原则", font_size=42, color=PAPER, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=0)

    principles = [
        ("一服务一数据库", "严格隔离，禁止跨库 Join，数据自治"),
        ("Nacos 配置中心", "Namespace 环境隔离，Group 分 DEV / PROD，动态刷新"),
        ("事件驱动优先", "MQ 解耦异步通信，同步调用保持扁平，禁止循环依赖"),
        ("接口幂等 + 分布式锁", "所有接口必须幂等，关键操作加分布式锁保护"),
        ("订单状态机", "全生命周期使用状态机管理，Spring Statemachine 驱动"),
        ("无状态 · 可降级 · 可观测", "服务必须无状态部署，支持降级，全链路可观测"),
    ]

    card_w = Inches(5.2)
    card_h = Inches(1.3)
    cols = 2
    gap_x = Inches(0.5)
    gap_y = Inches(0.3)
    start_x = Inches(0.8)
    start_y = Inches(3.2)

    for i, (title, desc) in enumerate(principles):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = INK_TINT
        shape.line.color.rgb = RGBColor(0x40, 0x40, 0x42)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.5), card_h - Inches(0.3))
        add_para(tf, title, font_size=14, color=PAPER, bold=True,
                 font_name="Microsoft YaHei UI", spacing_after=4)
        add_para(tf, desc, font_size=12, color=PAPER_TINT, spacing_after=0)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "架构设计原则", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "04 · Architecture", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 6: Microservices (light)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PAPER)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "微服务 · Services", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "05 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.5))
    add_para(tf, "MICROSERVICES · 服务拆分", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "8 大核心服务", font_size=42, color=INK, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=0)

    services = [
        ("gateway", "API 网关 — 统一入口，鉴权、限流、路由、跨域"),
        ("user-service", "用户服务 — 注册登录、JWT 鉴权、RBAC 权限体系"),
        ("merchant-service", "商家服务 — 商家入驻、菜品管理、购物车、SKU 规格"),
        ("order-service", "订单服务（最复杂）— 状态机、分布式事务、高并发优化"),
        ("payment-service", "支付服务 — 支付宝/微信沙箱集成、回调验签、退款"),
        ("delivery-service", "配送服务（最前沿）— 智能调度、实时位置、抢单/派单"),
        ("search-service", "搜索服务 — Elasticsearch 全文搜索与聚合"),
        ("notification-service", "通知服务 — WebSocket 推送、短信、消息可靠性保障"),
    ]

    tf = add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(4))

    for name, desc in services:
        add_para(tf, "─" * 80, font_size=6, color=PAPER_TINT,
                 font_name="Consolas", spacing_before=4, spacing_after=1)
        add_mixed_para(tf, [
            (f"  {name}", 13, INK, True, "Consolas"),
            (f"    {desc}", 13, INK_TINT, False, "Microsoft YaHei UI"),
        ], spacing_after=2)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "微服务拆分", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "05 · Services", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 7: Team (dark)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "团队 · Team", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "06 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.5))
    add_para(tf, "TEAM · 团队分工", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "5 人协作", font_size=42, color=PAPER, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=0)

    team = [
        ("成员 1 · 架构负责人", "全局把控", "网关服务、用户服务、Nacos 集群、全局框架、监控部署、代码审查"),
        ("成员 2", "商家服务", "商家入驻、菜品管理、购物车（Redis）、商家端订单处理"),
        ("成员 3 · 最复杂模块", "订单服务", "全流程状态机、Seata 分布式事务、高并发优化、事件驱动"),
        ("成员 4", "支付 + 通知", "支付沙箱集成、实时消息推送、RabbitMQ 死信队列、回调安全"),
        ("成员 5 · 最前沿模块", "配送调度", "骑手服务、智能调度算法（Redis ZSet + 贪心）、实时位置上报"),
    ]

    card_w = Inches(5.2)
    card_h = Inches(1.3)
    cols = 2
    gap_x = Inches(0.5)
    gap_y = Inches(0.25)
    start_x = Inches(0.8)
    start_y = Inches(3.0)

    for i, (role, name, desc) in enumerate(team):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = INK_TINT
        shape.line.color.rgb = RGBColor(0x40, 0x40, 0x42)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.25), y + Inches(0.12), card_w - Inches(0.5), card_h - Inches(0.2))
        add_para(tf, role.upper(), font_size=9, color=MUTED, font_name="Consolas", spacing_after=2)
        add_para(tf, name, font_size=16, color=PAPER, bold=True,
                 font_name="Microsoft YaHei UI", spacing_after=2)
        add_para(tf, desc, font_size=11, color=PAPER_TINT, spacing_after=0)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "5 人后端团队", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "06 · Team", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 8: Deployment (light)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PAPER)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "部署 · Deployment", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "07 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.5))
    add_para(tf, "DEPLOYMENT · 部署流程", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "Docker Compose 一键部署", font_size=42, color=INK, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=0)

    # Pipeline 1: Requirements
    tf = add_textbox(slide, Inches(0.8), Inches(2.8), Inches(11), Inches(0.3))
    add_para(tf, "环境要求 · REQUIREMENTS", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=0)

    req_items = [
        ("01", "JDK 21", "最新 LTS 版本"),
        ("02", "Maven 3.9+", "构建所有 JAR 包"),
        ("03", "Node.js 18+", "前端构建"),
        ("04", "Docker", "容器化部署"),
    ]

    step_w = Inches(2.6)
    step_h = Inches(1.4)
    step_y = Inches(3.3)
    step_gap = Inches(0.35)
    step_start = Inches(0.8)

    for i, (nb, title, desc) in enumerate(req_items):
        x = step_start + i * (step_w + step_gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, step_y, step_w, step_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PAPER_TINT
        shape.line.color.rgb = RGBColor(0xD0, 0xCE, 0xC8)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.2), step_y + Inches(0.15), step_w - Inches(0.4), step_h - Inches(0.2))
        add_para(tf, nb, font_size=14, color=MUTED, font_name="Playfair Display", spacing_after=2)
        add_para(tf, title, font_size=15, color=INK, bold=True,
                 font_name="Microsoft YaHei UI", spacing_after=2)
        add_para(tf, desc, font_size=11, color=INK_TINT, spacing_after=0)

    # Pipeline 2: Methods
    tf = add_textbox(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.3))
    add_para(tf, "部署方式 · METHODS", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=0)

    methods = [
        ("A", "本地开发", "--spring.profiles.active=local，连接 localhost"),
        ("B", "Docker Compose", "10 个容器一键启动，含监控"),
        ("C", "离线部署", "docker save/load 导出导入镜像"),
    ]

    method_w = Inches(3.5)
    method_h = Inches(1.2)
    method_y = Inches(5.4)
    method_gap = Inches(0.5)

    for i, (nb, title, desc) in enumerate(methods):
        x = step_start + i * (method_w + method_gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, method_y, method_w, method_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = PAPER_TINT
        shape.line.color.rgb = RGBColor(0xD0, 0xCE, 0xC8)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.2), method_y + Inches(0.12), method_w - Inches(0.4), method_h - Inches(0.2))
        add_para(tf, nb, font_size=14, color=MUTED, font_name="Playfair Display", spacing_after=2)
        add_para(tf, title, font_size=14, color=INK, bold=True,
                 font_name="Microsoft YaHei UI", spacing_after=2)
        add_para(tf, desc, font_size=11, color=INK_TINT, spacing_after=0)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "前端 :80 · API 网关 :9999 · Nacos :8848", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "07 · Deployment", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 9: Workflow (dark)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, INK)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "流程 · Workflow", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "08 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.5))
    add_para(tf, "WORKFLOW · 开发协作", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=8)
    add_para(tf, "从分支到上线", font_size=42, color=PAPER, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=0)

    workflow = [
        ("分支策略", "dev 分支 → feature/xxx 分支开发 → PR → 架构负责人 Review → 合并"),
        ("跨成员协作", "数据库整体设计（雪花 ID）、统一接口规范、Knife4J 文档、ER 关系对齐"),
        ("测试体系", "JUnit5 单元测试 + Testcontainers 集成测试 + JMeter 性能压测（订单创建、骑手抢单）"),
        ("可观测性", "traceId 全链路追踪 · Grafana 面板（订单量、骑手分布、事务成功率）"),
    ]

    card_w = Inches(5.2)
    card_h = Inches(1.4)
    cols = 2
    gap_x = Inches(0.5)
    gap_y = Inches(0.3)
    start_x = Inches(0.8)
    start_y = Inches(3.2)

    for i, (title, desc) in enumerate(workflow):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = INK_TINT
        shape.line.color.rgb = RGBColor(0x40, 0x40, 0x42)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

        tf = add_textbox(slide, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.5), card_h - Inches(0.3))
        add_para(tf, title, font_size=14, color=PAPER, bold=True,
                 font_name="Microsoft YaHei UI", spacing_after=6)
        add_para(tf, desc, font_size=12, color=PAPER_TINT, spacing_after=0)

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "开发流程与协作", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "08 · Workflow", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════════════════════
    # Page 10: Finale (hero light)
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, PAPER)

    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.3))
    add_para(tf, "结语 · Finale", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(0.4), Inches(2), Inches(0.3))
    add_para(tf, "10 / 10", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(4.5))

    add_para(tf, "TAKEAWAY · 核心关键词", font_size=10, color=MUTED,
             font_name="Consolas", spacing_after=16)

    add_para(tf, "云原生外卖平台", font_size=48, color=INK, bold=True,
             font_name="Microsoft YaHei UI", spacing_after=24)

    add_mixed_para(tf, [
        ("云原生", 14, MUTED, False, "Consolas"),
        ("  ·  ", 14, MUTED, False, "Consolas"),
        ("微服务", 14, MUTED, False, "Consolas"),
        ("  ·  ", 14, MUTED, False, "Consolas"),
        ("事件驱动", 14, MUTED, False, "Consolas"),
        ("  ·  ", 14, MUTED, False, "Consolas"),
        ("实时智能", 14, MUTED, False, "Consolas"),
        ("  ·  ", 14, MUTED, False, "Consolas"),
        ("全链路可观测", 14, MUTED, False, "Consolas"),
    ], spacing_after=24)

    add_para(tf, "从架构设计到容器部署，从状态机到智能调度，这是一个完整的微服务实践。",
             font_size=16, color=INK_TINT, spacing_after=32)

    add_mixed_para(tf, [
        ("SpringBoot 外卖点餐平台", 11, MUTED, False, "Consolas"),
        ("  ·  ", 11, MUTED, False, "Consolas"),
        ("2026", 11, MUTED, False, "Consolas"),
    ])

    # Footer
    tf = add_textbox(slide, Inches(0.8), Inches(6.8), Inches(6), Inches(0.3))
    add_para(tf, "感谢聆听", font_size=9, color=MUTED, font_name="Consolas")
    tf = add_textbox(slide, Inches(10.5), Inches(6.8), Inches(2), Inches(0.3))
    add_para(tf, "— · —", font_size=9, color=MUTED, font_name="Consolas", alignment=PP_ALIGN.RIGHT)

    # ── Save ──
    out_path = os.path.join(os.path.dirname(__file__), "SpringBoot外卖点餐平台.pptx")
    prs.save(out_path)
    print(f"PPTX saved to: {out_path}")


if __name__ == "__main__":
    create_pptx()
